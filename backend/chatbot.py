import google.generativeai as genai
import os
import io
import uuid
import docx
import PyPDF2
from pptx import Presentation
import json
import re
from dotenv import load_dotenv 
from sentence_transformers import SentenceTransformer, util

# Configure the Gemini API

# Initialize the Gemini model
model = genai.GenerativeModel("gemini-2.0-flash")

def generate_chat_response(user_message):
    """
    Sends a user message to the Gemini API and returns the AI-generated response.
    """
    print("genrate chat response")
    prompt = f"""You are an AI assistant specialized in helping students prepare for UPSC and MPSC exams.
        Your goal is to provide accurate, helpful information about exam preparation, syllabus, 
        study materials, strategies, and answer general questions.
        
        User question: {user_message}
        """
        
    try:
        response = model.generate_content(
            prompt,
            generation_config=genai.types.GenerationConfig(
                candidate_count=1,
                max_output_tokens=400,  # Adjust based on expected response length
                temperature=0.7
            )
        )
        return response.text.strip()
    except Exception as e:
        return f"Error: {str(e)}"


def extract_text_from_file(file, file_extension):
    file_content = file.read()
    
    if file_extension == 'pdf':
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
        page_texts = []
        for i, page in enumerate(pdf_reader.pages):
            text = page.extract_text()
            if text:
                page_texts.append({'page': i+1, 'text': text})
        return page_texts

    elif file_extension == 'txt':
        return [{'page': 1, 'text': file_content.decode('utf-8', errors='replace')}]

    elif file_extension == 'docx':
        doc = docx.Document(io.BytesIO(file_content))
        text = '\n'.join([para.text for para in doc.paragraphs])
        return [{'page': 1, 'text': text}]

    elif file_extension == 'pptx':
        ppt = Presentation(io.BytesIO(file_content))
        text = '\n'.join(
            shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")
        )
        return [{'page': 1, 'text': text}]

    return []

def chunk_text_by_page(page_texts, chunk_size=500):
    chunks = []
    for page_obj in page_texts:
        page = page_obj['page']
        text = page_obj['text']
        sentences = text.split('. ')
        chunk = ""
        for sentence in sentences:
            if len(chunk.split()) + len(sentence.split()) <= chunk_size:
                chunk += sentence + ". "
            else:
                chunks.append({'text': chunk.strip(), 'page': page})
                chunk = sentence + ". "
        if chunk:
            chunks.append({'text': chunk.strip(), 'page': page})
    return chunks

def retrieve_relevant_chunks(chunks, query, top_k=3):
    model = SentenceTransformer('all-MiniLM-L6-v2')
    texts = [c['text'] for c in chunks]
    chunk_embeddings = model.encode(texts, convert_to_tensor=True)
    query_embedding = model.encode(query, convert_to_tensor=True)
    hits = util.semantic_search(query_embedding, chunk_embeddings, top_k=top_k)
    return [chunks[hit['corpus_id']] for hit in hits[0]]

def save_uploaded_document(file, upload_folder):
    if file.filename == '':
        return None, "No selected file"
    
    file_extension = file.filename.split('.')[-1].lower()
    if file_extension not in ['pdf', 'docx', 'pptx', 'txt']:
        return None, "Unsupported file format"
    
    try:
        document_id = str(uuid.uuid4())
        file_path = os.path.join(upload_folder, document_id + '.' + file_extension)
        file.seek(0)
        file.save(file_path)

        file.seek(0)
        page_texts = extract_text_from_file(file, file_extension)
        
        # Store the content for future reference
        content = []
        for page in page_texts:
            content.append(page['text'])
        
        return {
            'id': document_id,
            'filename': file.filename,
            'path': file_path,
            'page_texts': page_texts,
            'content': '\n'.join(content)
        }, None

    except Exception as e:
        return None, str(e)

def generate_document_response(document_path, file_extension, user_message, model_name="gemini-2.0-flash"):
    """
    Generate a response with document reference highlights
    """
    try:
        page_contexts = []
        
        if file_extension == 'pdf':
            # Read PDF content page by page
            pdf_reader = PyPDF2.PdfReader(document_path)
            for i, page in enumerate(pdf_reader.pages):
                text = page.extract_text()
                if text:
                    page_contexts.append({'page': i + 1, 'text': text.strip()})
        elif file_extension == 'txt':
            with open(document_path, 'r', encoding='utf-8', errors='replace') as f:
                text = f.read()
                page_contexts.append({'page': 1, 'text': text.strip()})
        elif file_extension == 'docx':
            doc = docx.Document(document_path)
            text = '\n'.join([para.text for para in doc.paragraphs])
            page_contexts.append({'page': 1, 'text': text.strip()})
        elif file_extension == 'pptx':
            ppt = Presentation(document_path)
            text = '\n'.join(
                shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")
            )
            page_contexts.append({'page': 1, 'text': text.strip()})

        # Build prompt with paginated content (limit tokens)
        limited_context = ""
        included_pages = []
        for ctx in page_contexts:
            if len(limited_context) + len(ctx['text']) > 8000:
                break
            limited_context += f"\n\n(Page {ctx['page']}):\n{ctx['text']}"
            included_pages.append(ctx)

        prompt = f"""
        You are a helpful assistant. The following is content from a document split by pages.
        Based on this content, answer the user query at the end.

        {limited_context}

        User Query: {user_message}

        If the answer cannot be found in the document, say so.
        Also indicate which parts of the document (by page and a few words) were referred to.
        """

        # Get Gemini response
        model = genai.GenerativeModel(model_name)
        response = model.generate_content(prompt)
        answer = response.text

        # Simple keyword-based match to extract referenced chunks (demo logic)
        import difflib
        matched_references = []
        for page_ctx in included_pages:
            matches = difflib.get_close_matches(user_message, [page_ctx['text']], n=1, cutoff=0.3)
            if matches:
                matched_references.append({
                    'page': page_ctx['page'],
                    'text': matches[0][:300]  # limit snippet length
                })

        return answer, matched_references, None
    except Exception as e:
        return None, None, str(e)

    
# def genrate_rag_response(question , context , links):
#     """
#     Sends a user message to the Gemini API and returns the AI-generated response.
#     """
    
#     prompt = f"""You are best proffesor for mpsc and upsc an you are good at answering form a context and adding a more infomatice contetn to answer.
#                   I will provide you a context from wich you give the answerr also add your knowledge to the answer related to question but take most of from context.
#                   In response also add the link that i provide to you and also add your referse with the answer.

#                   Context is {context}
#                   And user question is {question} give me answer in json format with keys as answer and reference.
#                   Also add the link that i provide {links} to you in the answer.
                  
        
#         """
        
#     try:
#         response = model.generate_content(
#             prompt,
#             generation_config=genai.types.GenerationConfig(
#                 candidate_count=1,
#                 max_output_tokens=400,  # Adjust based on expected response length
#                 temperature=0.7
#             )
#         )
#         return response.text.strip()
#     except Exception as e:
#         return f"Error: {str(e)}"

# def genrate_rag_response(question, context, links):
#     """
#     Sends a user message to the Gemini API and returns the AI-generated response
#     in JSON format with 'answer' and 'reference' keys.
#     """
    
#     # Clean and format the links as a list of strings
#     links_list = links if links else ["No reference provided"]

#     # Strict JSON prompt for Gemini to follow
#     prompt = f"""
#     You are an expert professor specializing in MPSC and UPSC exam preparation. Your task is to answer user questions by referring to the given context and supplementing the answer with relevant knowledge.
    
#     ### Instructions:
#     - what ever asked provide the answer in json format with keys as answer and reference.
#     - Use the provided context to generate a detailed and well-structured response.
#     - Use most of the context to answer the question.
#     - Ensure the response is informative and relevant to the question asked.
#     - If the answer cannot be found in the context, provide a relevant and accurate response.
#     - Add any useful information that complements the context.
#     - Include the provided link(s) as references in the response.
#     - Strictly return the result in the following JSON format:
    
#     ```json
#     {{
#         "answer": "Your answer here, based on the context and your knowledge.",
#         "reference": {links_list}
#     }}
#     ```

#     ### Provided Context:
#     {context}
    
#     ### User Question:
#     {question}
    
#     ### Provided Links:
#     {", ".join(links_list)}
#     """

#     try:
#         # Send the prompt to the Gemini API
#         response = model.generate_content(
#             prompt,
#             generation_config=genai.types.GenerationConfig(
#                 candidate_count=1,
#                 max_output_tokens=400,
#                 temperature=0.7
#             )
#         )

#         # Extract and clean raw response
#         raw_response = response.text.strip()
#         #print("Raw response from AI:", raw_response)

#         # ✅ Handle Markdown if it wraps the JSON with ```json or ```
#         if raw_response.startswith("```json"):
#             raw_response = raw_response[7:-3].strip()  # Remove ```json and ending ```
#         elif raw_response.startswith("```"):
#             raw_response = raw_response[3:-3].strip()  # Remove generic markdown ```

#         # ✅ Try parsing the response as JSON
#         try:
#             response_json = json.loads(raw_response)
#             #print("Parsed JSON response:", response_json)
            
#             # Ensure required keys are present
#             if "answer" in response_json and "reference" in response_json:
#                 return response_json
#             else:
#                 return {
#                     "error": "Response does not contain the required keys.",
#                     "raw_response": raw_response
#                 }
#         except json.JSONDecodeError:
#             return {
#                 "error": "Invalid JSON response from AI.",
#                 "raw_response": raw_response
#             }

#     except Exception as e:
#         return {
#             "error": str(e)
#         }



def generate_rag_response(question, context, links):
    """
    Sends a user message to the Gemini API and returns the AI-generated response
    in JSON format with 'answer' and 'reference' keys.
    """

    # Clean and format the links as a list of strings
    links_list = links if links else ["No reference provided"]

    # Strict JSON prompt for Gemini to follow
    prompt = f"""
    You are an expert professor specializing in MPSC and UPSC exam preparation. Your task is to answer user questions by referring to the given context and supplementing the answer with relevant knowledge.

    ### Instructions:
    - Provide the answer strictly in JSON format with these keys:
      - "answer" - A detailed, informative, and accurate response based on the context provided.
      - "reference" - A list of reference links (from the provided links).
    - Strictly return the JSON format as shown below:
    
    Example 1:
    ```json
    {{
        "answer": "The capital of India is New Delhi. It was officially declared the capital in 1911.",
        "reference": ["https://example.com/india-capital"]
    }}
    ```

    Example 2:
    ```json
    {{
        "answer": "The Civil Services Examination is conducted by UPSC in three stages: Preliminary, Mains, and Interview.",
        "reference": ["https://upsc.gov.in"]
    }}
    ```

    - If the answer cannot be found in the context, provide a well-informed response.
    - Use most of the context provided to answer the question.
    - Return the JSON as the only response without any additional comments or explanations.

    ### Provided Context:
    {context}
    
    ### User Question:
    {question}
    
    ### Provided Links:
    {", ".join(links_list)}
    """

    try:
        # Send the prompt to the Gemini API
        response = model.generate_content(
            prompt,
            generation_config=genai.types.GenerationConfig(
                candidate_count=1,
                max_output_tokens=500,  # Increase to prevent truncation
                temperature=0.3  # Lower temp for consistent responses
            )
        )

        # Extract and clean raw response
        raw_response = response.text.strip()

        # ✅ Use regex to clean up any unwanted markdown or text
        raw_response = re.sub(r"```json|```", "", raw_response).strip()

        # ✅ Try parsing the response as JSON
        try:
            response_json = json.loads(raw_response)

            # Ensure required keys are present
            if "answer" in response_json and "reference" in response_json:
                return response_json
            else:
                return {
                    "error": "Response does not contain required keys.",
                    "raw_response": raw_response
                }

        except json.JSONDecodeError:
            # If JSON parsing fails, try to extract the possible JSON using regex
            possible_json_match = re.search(r"{.*}", raw_response, re.DOTALL)

            if possible_json_match:
                try:
                    # Try parsing the extracted JSON
                    cleaned_json = json.loads(possible_json_match.group())
                    if "answer" in cleaned_json and "reference" in cleaned_json:
                        return cleaned_json
                except json.JSONDecodeError:
                    pass

            # Return error if everything fails
            return {
                "error": "Invalid JSON response from AI.",
                "raw_response": raw_response
            }

    except Exception as e:
        # Catch any other errors
        return {
            "error": str(e)
        }
